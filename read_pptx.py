#!/usr/bin/env python3
"""Structured PPTX reader — extracts slides, notes, tables, and image metadata.

Outputs structured markdown (default) or JSON for programmatic use.

Usage:
    ./venv/bin/python read_pptx.py <file.pptx>
    ./venv/bin/python read_pptx.py <file.pptx> --slides 1,3,5
    ./venv/bin/python read_pptx.py <file.pptx> --json
    ./venv/bin/python read_pptx.py <file.pptx> --notes-only
    ./venv/bin/python read_pptx.py <file.pptx> --summary
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu


def emu_to_inches(emu_val):
    """Convert EMU to inches, rounded to 2 decimals."""
    if emu_val is None:
        return None
    return round(emu_val / 914400, 2)


def extract_table(shape):
    """Extract table data as list of rows (list of cell strings)."""
    table = shape.table
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    return rows


def table_to_markdown(rows):
    """Convert a list of rows to a markdown table."""
    if not rows:
        return ""
    # Use first row as header
    header = rows[0]
    col_widths = [max(len(str(cell)) for cell in col) for col in zip(*rows)]
    col_widths = [max(w, 3) for w in col_widths]

    lines = []
    # Header
    line = "| " + " | ".join(str(h).ljust(w) for h, w in zip(header, col_widths)) + " |"
    lines.append(line)
    # Separator
    line = "| " + " | ".join("-" * w for w in col_widths) + " |"
    lines.append(line)
    # Data rows
    for row in rows[1:]:
        line = "| " + " | ".join(str(c).ljust(w) for c, w in zip(row, col_widths)) + " |"
        lines.append(line)
    return "\n".join(lines)


def extract_slide(slide, slide_number, prs):
    """Extract structured data from a single slide."""
    data = {
        "number": slide_number,
        "layout": None,
        "title": None,
        "text_blocks": [],
        "tables": [],
        "images": [],
        "charts": [],
        "shapes": [],
        "notes": None,
    }

    # Layout name
    try:
        data["layout"] = slide.slide_layout.name
    except Exception:
        pass

    # Title
    if slide.shapes.title:
        data["title"] = slide.shapes.title.text.strip()

    # Process all shapes
    for shape in slide.shapes:
        pos = {
            "left": emu_to_inches(shape.left),
            "top": emu_to_inches(shape.top),
            "width": emu_to_inches(shape.width),
            "height": emu_to_inches(shape.height),
        }

        # Tables
        if shape.has_table:
            rows = extract_table(shape)
            data["tables"].append({"rows": rows, "position": pos})
            continue

        # Charts
        if shape.has_chart:
            chart = shape.chart
            chart_data = {
                "type": str(chart.chart_type),
                "title": chart.chart_title.text_frame.text if chart.has_title else None,
                "position": pos,
            }
            # Try to extract series names
            try:
                chart_data["series"] = [s.name for s in chart.series if s.name]
            except Exception:
                pass
            data["charts"].append(chart_data)
            continue

        # Images
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            img_info = {"position": pos, "name": shape.name}
            try:
                img_info["content_type"] = shape.image.content_type
                img_info["size_bytes"] = len(shape.image.blob)
                img_info["filename"] = shape.image.filename
            except Exception:
                pass
            data["images"].append(img_info)
            continue

        # Text content (excluding title, already captured)
        if shape.has_text_frame and shape.text.strip():
            if shape == slide.shapes.title:
                continue
            paragraphs = []
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    p_info = {"text": text}
                    # Detect formatting
                    if para.runs:
                        run = para.runs[0]
                        if run.font.bold:
                            p_info["bold"] = True
                        if run.font.size:
                            p_info["font_size"] = round(run.font.size.pt, 1)
                    paragraphs.append(p_info)
            if paragraphs:
                data["text_blocks"].append({
                    "paragraphs": paragraphs,
                    "position": pos,
                    "shape_name": shape.name,
                })

    # Speaker notes
    try:
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                data["notes"] = notes_text
    except Exception:
        pass

    return data


def extract_presentation(pptx_path, slide_numbers=None):
    """Extract structured data from an entire presentation."""
    prs = Presentation(pptx_path)

    result = {
        "file": str(pptx_path),
        "slide_width": emu_to_inches(prs.slide_width),
        "slide_height": emu_to_inches(prs.slide_height),
        "total_slides": len(prs.slides),
        "slides": [],
    }

    # Core properties
    try:
        cp = prs.core_properties
        props = {}
        if cp.title:
            props["title"] = cp.title
        if cp.author:
            props["author"] = cp.author
        if cp.subject:
            props["subject"] = cp.subject
        if cp.created:
            props["created"] = cp.created.isoformat()
        if cp.modified:
            props["modified"] = cp.modified.isoformat()
        if props:
            result["properties"] = props
    except Exception:
        pass

    for i, slide in enumerate(prs.slides, 1):
        if slide_numbers and i not in slide_numbers:
            continue
        result["slides"].append(extract_slide(slide, i, prs))

    return result


def format_markdown(data, notes_only=False, summary=False):
    """Format extracted data as structured markdown."""
    lines = []

    if summary:
        lines.append(f"# {data.get('properties', {}).get('title', Path(data['file']).stem)}")
        lines.append("")
        lines.append(f"**File:** `{data['file']}`")
        lines.append(f"**Slides:** {data['total_slides']} | **Size:** {data['slide_width']}\" x {data['slide_height']}\"")
        props = data.get("properties", {})
        if props.get("author"):
            lines.append(f"**Author:** {props['author']}")
        if props.get("modified"):
            lines.append(f"**Modified:** {props['modified']}")
        lines.append("")
        lines.append("## Slide Overview")
        lines.append("")
        for slide in data["slides"]:
            title = slide["title"] or "(no title)"
            extras = []
            if slide["tables"]:
                extras.append(f"{len(slide['tables'])} table(s)")
            if slide["images"]:
                extras.append(f"{len(slide['images'])} image(s)")
            if slide["charts"]:
                extras.append(f"{len(slide['charts'])} chart(s)")
            if slide["notes"]:
                extras.append("has notes")
            suffix = f" — {', '.join(extras)}" if extras else ""
            lines.append(f"{slide['number']}. **{title}**{suffix}")
        return "\n".join(lines)

    # Full output
    lines.append(f"# {data.get('properties', {}).get('title', Path(data['file']).stem)}")
    lines.append("")
    lines.append(f"**File:** `{data['file']}`  ")
    lines.append(f"**Slides:** {data['total_slides']} | **Dimensions:** {data['slide_width']}\" x {data['slide_height']}\"")
    props = data.get("properties", {})
    if props.get("author"):
        lines.append(f"**Author:** {props['author']}")
    if props.get("modified"):
        lines.append(f"**Modified:** {props['modified']}")
    lines.append("")

    for slide in data["slides"]:
        lines.append(f"---")
        lines.append(f"## Slide {slide['number']}: {slide['title'] or '(untitled)'}")
        if slide["layout"]:
            lines.append(f"*Layout: {slide['layout']}*")
        lines.append("")

        if not notes_only:
            # Text blocks
            for block in slide["text_blocks"]:
                for para in block["paragraphs"]:
                    text = para["text"]
                    if para.get("bold"):
                        text = f"**{text}**"
                    lines.append(text)
                lines.append("")

            # Tables
            for table in slide["tables"]:
                lines.append(table_to_markdown(table["rows"]))
                lines.append("")

            # Charts
            for chart in slide["charts"]:
                chart_title = chart.get("title") or "Untitled"
                lines.append(f"[Chart: {chart['type']} — \"{chart_title}\"]")
                if chart.get("series"):
                    lines.append(f"  Series: {', '.join(chart['series'])}")
                lines.append("")

            # Images
            for img in slide["images"]:
                name = img.get("filename") or img.get("name", "image")
                size = img.get("size_bytes")
                size_str = f" ({size // 1024}KB)" if size else ""
                ctype = img.get("content_type", "")
                lines.append(f"[Image: {name}{size_str} {ctype}]")
            if slide["images"]:
                lines.append("")

        # Notes
        if slide["notes"]:
            lines.append(f"> **Notes:** {slide['notes']}")
            lines.append("")

    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(
        description="Extract structured content from PPTX files"
    )
    parser.add_argument("file", help="Path to .pptx file")
    parser.add_argument(
        "--slides",
        help="Comma-separated slide numbers to extract (e.g., 1,3,5)",
    )
    parser.add_argument(
        "--json",
        action="store_true",
        dest="output_json",
        help="Output as JSON instead of markdown",
    )
    parser.add_argument(
        "--notes-only",
        action="store_true",
        help="Only extract speaker notes",
    )
    parser.add_argument(
        "--summary",
        action="store_true",
        help="One-line-per-slide overview",
    )

    args = parser.parse_args()

    pptx_path = Path(args.file)
    if not pptx_path.exists():
        print(f"Error: {pptx_path} does not exist", file=sys.stderr)
        sys.exit(1)

    slide_numbers = None
    if args.slides:
        slide_numbers = [int(s.strip()) for s in args.slides.split(",")]

    data = extract_presentation(pptx_path, slide_numbers)

    if args.output_json:
        print(json.dumps(data, indent=2, ensure_ascii=False))
    else:
        print(format_markdown(data, notes_only=args.notes_only, summary=args.summary))


if __name__ == "__main__":
    main()
